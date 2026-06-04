"""
Build the Voice Coaching System PowerPoint presentation.
Based entirely on VocalCoach_Kim source code — not the framework PDF.

Run:  python3 build_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1A, 0x3A, 0x5C)
MID_BLUE   = RGBColor(0x1F, 0x77, 0xB4)
LIGHT_BLUE = RGBColor(0xAE, 0xC7, 0xE8)
ACCENT     = RGBColor(0xFF, 0x7F, 0x0E)
GREEN      = RGBColor(0x2C, 0xA0, 0x2C)
PURPLE     = RGBColor(0x6A, 0x3D, 0x9A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK  = RGBColor(0x1A, 0x1A, 0x1A)
LGREY      = RGBColor(0xF5, 0xF5, 0xF5)
MGREY      = RGBColor(0xDD, 0xDD, 0xDD)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Low-level helpers ─────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def box(slide, x, y, w, h, fill=None, line=None):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
    else:
        sh.line.fill.background()
    return sh


def txt(slide, text, x, y, w, h,
        size=16, bold=False, italic=False,
        color=TEXT_DARK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def header(slide, title, sub=None):
    box(slide, 0, 0, 13.33, 1.05, fill=DARK_BLUE)
    txt(slide, title, 0.3, 0.08, 12.5, 0.6, size=26, bold=True, color=WHITE)
    if sub:
        txt(slide, sub, 0.3, 0.65, 12.5, 0.34, size=13, color=LIGHT_BLUE)
    box(slide, 0, 1.05, 13.33, 0.04, fill=ACCENT)


def formula(slide, text, x, y, w, h=0.44):
    box(slide, x, y, w, h, fill=RGBColor(0xE8, 0xF0, 0xFE),
        line=MID_BLUE)
    txt(slide, text, x+0.1, y+0.04, w-0.2, h-0.08,
        size=13, italic=True, color=DARK_BLUE)


def table(slide, headers, rows, x, y, widths,
          row_h=0.38, hdr_bg=LIGHT_BLUE, alt=LGREY, fsize=12):
    def draw_row(data, ry, bg, bold=False):
        cx = x
        for cell, cw in zip(data, widths):
            box(slide, cx, ry, cw, row_h, fill=bg, line=MGREY)
            txt(slide, str(cell), cx+0.06, ry+0.04, cw-0.12, row_h-0.08,
                size=fsize, bold=bold, color=DARK_BLUE if bold else TEXT_DARK, wrap=True)
            cx += cw
    draw_row(headers, y, hdr_bg, bold=True)
    for i, row in enumerate(rows):
        draw_row(row, y + row_h*(i+1), alt if i % 2 == 0 else WHITE)


def blist(slide, items, x, y, w, fsize=13, gap=0.38, color=TEXT_DARK):
    for i, item in enumerate(items):
        txt(slide, "  •  " + item, x, y + i*gap, w, gap,
            size=fsize, color=color)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    sl = add_slide(prs)
    box(sl, 0, 0, 13.33, 7.5, fill=DARK_BLUE)
    box(sl, 0, 4.7, 13.33, 2.8, fill=MID_BLUE)
    box(sl, 0, 4.65, 13.33, 0.08, fill=ACCENT)

    txt(sl, "Voice Coaching System",
        1.0, 1.1, 11.33, 1.15, size=44, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "for Singing Performance Evaluation",
        1.0, 2.2, 11.33, 0.7, size=30, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    txt(sl, "Reference-guided pipeline  ·  "
           "Three-model architecture  ·  Deterministic scoring",
        1.0, 3.1, 11.33, 0.5, size=16,
        color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.CENTER)
    txt(sl, "GTSinger Dataset  ·  MusicXML + TextGrid  ·  "
           "PyTorch / Wav2Vec2 / torchcrepe",
        1.0, 4.85, 11.33, 0.45, size=14, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "IMAPLE Lab  ·  Drexel University  ·  2026",
        1.0, 5.6, 11.33, 0.4, size=15, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)


def slide_overview(prs):
    sl = add_slide(prs)
    header(sl, "System Overview",
           "Audio in → feature extraction → compare to reference → score → interpret")

    stages = [
        ("1", "WAV Audio",
         "16 kHz mono singing recording", DARK_BLUE),
        ("2", "Feature Extraction",
         "Pitch (F0 + voiced mask)\nNote onset/offset times\nPhoneme boundaries",
         MID_BLUE),
        ("3", "Reference Data",
         "MusicXML → note pitch, onset, offset, tempo\nTextGrid → phoneme & word boundaries",
         RGBColor(0x27, 0x6F, 0x96)),
        ("4", "Metric Computation",
         "Pitch error · Timing error\nDuration error · Phoneme error",
         ACCENT),
        ("5", "Scoring Engine",
         "Piecewise normalise → Category scores\n→ Final score [0–100]",
         GREEN),
        ("6", "Interpretation",
         "Strengths  ·  Weaknesses\nRule-based coaching feedback",
         PURPLE),
    ]
    for i, (num, title, body, clr) in enumerate(stages):
        col, row = i // 3, i % 3
        bx, by = 0.3 + col*6.55, 1.2 + row*2.05
        box(sl, bx, by, 6.1, 1.9, fill=clr)
        txt(sl, num,   bx+0.1, by+0.06, 0.5, 0.45, size=22, bold=True, color=WHITE)
        txt(sl, title, bx+0.55, by+0.08, 5.4, 0.42, size=15, bold=True, color=WHITE)
        txt(sl, body,  bx+0.55, by+0.54, 5.4, 1.26, size=12, color=WHITE)

    txt(sl, "▶", 6.32, 3.1, 0.6, 0.5, size=26, bold=True,
        color=DARK_BLUE, align=PP_ALIGN.CENTER)


def slide_parameters(prs):
    sl = add_slide(prs)
    header(sl, "Parameters We Track",
           "What the system extracts from audio and reads from reference files")

    # LEFT — performance side
    box(sl, 0.25, 1.15, 6.1, 6.2, fill=RGBColor(0xE8, 0xF1, 0xFB))
    txt(sl, "Performance Side  (extracted from WAV)",
        0.38, 1.22, 5.8, 0.38, size=14, bold=True, color=DARK_BLUE)

    perf = [
        ("f̂ₜ",     "Frame-level Fundamental Frequency (F0)",
         "Estimated sung pitch at every frame · units: Hz · 0 = unvoiced"),
        ("v̂ₜ",     "Voiced / Unvoiced Decision",
         "1 = singing frame · 0 = silence / breath\nGates pitch scoring — only voiced frames are evaluated"),
        ("ŝₙ",     "Note Onset Time",
         "Predicted start time of note n (seconds)"),
        ("êₙ",     "Note Offset Time",
         "Predicted end time of note n (seconds)"),
        ("âₚ, b̂ₚ", "Phoneme Boundaries",
         "Predicted start & end times for each phoneme p (seconds)"),
    ]
    for i, (sym, name, desc) in enumerate(perf):
        by = 1.66 + i*1.0
        box(sl, 0.38, by, 5.82, 0.9, fill=WHITE, line=MID_BLUE)
        txt(sl, sym,  0.46, by+0.04, 0.9, 0.42, size=17, bold=True, color=MID_BLUE)
        txt(sl, name, 1.32, by+0.04, 4.8, 0.34, size=12, bold=True, color=DARK_BLUE)
        txt(sl, desc, 1.32, by+0.42, 4.8, 0.42, size=10, color=TEXT_DARK)

    # RIGHT — reference side
    box(sl, 6.65, 1.15, 6.45, 6.2, fill=RGBColor(0xFD, 0xF3, 0xE3))
    txt(sl, "Reference Side  (from MusicXML + TextGrid)",
        6.78, 1.22, 6.2, 0.38, size=14, bold=True, color=DARK_BLUE)

    ref_groups = [
        ("MusicXML", ACCENT, [
            "mₙ  —  MIDI pitch of reference note n",
            "rₙ = 440 · 2^((mₙ−69)/12)  —  reference frequency (Hz)",
            "sₙ  —  Reference note onset time (seconds)",
            "eₙ  —  Reference note offset time (seconds)",
            "dₙ = eₙ − sₙ  —  Reference note duration",
            "tempo (BPM), time signature",
        ]),
        ("TextGrid", GREEN, [
            "aₚ  —  Reference phoneme start time",
            "bₚ  —  Reference phoneme end time",
            "Word-level boundary times for lyric scoring",
        ]),
    ]
    cy = 1.62
    for src, clr, items in ref_groups:
        box(sl, 6.78, cy, 6.2, 0.34, fill=clr)
        txt(sl, src, 6.86, cy+0.04, 6.0, 0.28, size=13, bold=True, color=WHITE)
        cy += 0.38
        for item in items:
            txt(sl, "  • " + item, 6.82, cy, 6.2, 0.33, size=11, color=TEXT_DARK)
            cy += 0.34
        cy += 0.1

    formula(sl,
            "Voiced set:  V  =  { t : f̂ₜ > 0  and  rₜ > 0 }  —  "
            "frames where both sung voice and reference are voiced",
            0.25, 7.12, 12.8, 0.38)


def slide_pitch_metrics(prs):
    sl = add_slide(prs)
    header(sl, "Metrics — Pitch",
           "Four pitch metrics computed from AlignmentResult · tolerance τ = 50¢ (config: metrics.pitch.cents_tolerance)")

    items = [
        ("Pitch Accuracy  (pitch_accuracy)",
         "pitch_accuracy  =  mean( |dev| ≤ τ )  over all matched notes with valid pitch",
         "Fraction of matched note pairs whose pitch deviation is within tolerance.\n"
         "dev = NoteAlignmentMatch.pitch_deviation_cents  (positive = predicted is sharper).\n"
         "Returns None when no matched note has pitch data.  Higher is better.",
         MID_BLUE),
        ("Mean Absolute Cent Error  (mace_cents)",
         "MACE  =  mean( |dev| )  over all matched notes",
         "Average absolute pitch deviation in cents — no tolerance threshold needed.\n"
         "Scoring breakpoints: 0¢→100,  25¢→88,  50¢→75,  100¢→50,  200¢→0.\n"
         "(50¢ ≈ quarter-tone;  100¢ = 1 semitone)   Lower is better.",
         RGBColor(0x27, 0x6F, 0x96)),
        ("Pitch RMSE  (pitch_rmse_cents)",
         "PitchRMSE  =  √( mean( dev² ) )",
         "Root-mean-square cent error — penalises large deviations more heavily than MACE.\n"
         "Scoring breakpoints: 0¢→100,  25¢→88,  50¢→72,  100¢→45,  200¢→0.\n"
         "More sensitive to occasional large pitch jumps.   Lower is better.",
         RGBColor(0x1A, 0x5C, 0x96)),
        ("Note-Level Pitch Accuracy  (note_pitch_accuracy)",
         "note_pitch_accuracy  =  accuracy   (set equal to pitch_accuracy in build_pitch_metrics())",
         "Same numeric value as pitch_accuracy — build_pitch_metrics() sets both to the same result.\n"
         "Per-note detail is stored separately in PitchMetrics.per_note (List[MetricBreakdown]).\n"
         "Coaching interpretation: '% of notes sung in tune.'   Higher is better.",
         ACCENT),
    ]

    for i, (title, form, desc, clr) in enumerate(items):
        col, row = i % 2, i // 2
        bx, by = 0.25 + col*6.55, 1.18 + row*3.05
        box(sl, bx, by, 6.2, 2.92, fill=RGBColor(0xF0, 0xF5, 0xFF), line=clr)
        box(sl, bx, by, 6.2, 0.36, fill=clr)
        txt(sl, title, bx+0.1, by+0.04, 6.0, 0.3, size=13, bold=True, color=WHITE)
        box(sl, bx+0.1, by+0.42, 6.0, 0.56, fill=RGBColor(0xE0, 0xEA, 0xF8), line=clr)
        txt(sl, form, bx+0.18, by+0.46, 5.84, 0.48, size=12, italic=True, color=DARK_BLUE)
        txt(sl, desc, bx+0.1, by+1.06, 6.0, 1.78, size=11, color=TEXT_DARK)


def slide_timing_metrics(prs):
    sl = add_slide(prs)
    header(sl, "Metrics — Timing",
           "Onset/offset error and IOI consistency · tolerance τ = 50 ms (config: metrics.timing.onset_tolerance_ms)")

    items = [
        ("Onset Error  (mean_onset_error_ms)",
         "mean_onset_error_ms  =  mean( onset_dev_s × 1000 )  in ms\n"
         "std_onset_error_ms   =  std( onset_dev_s × 1000 )\n"
         "mean_abs_onset_error_ms  =  mean( |onset_dev_s| × 1000 )",
         "Signed deviation: positive = predicted is LATER than reference.\n"
         "All three stats reported — mean (bias), std (consistency), MAE (accuracy).\n"
         "Scoring uses MAE breakpoints: 0→100, 25→88, 50→75, 100→50, 200→0 ms.",
         MID_BLUE),
        ("Offset Error  (mean_offset_error_ms)",
         "mean_offset_error_ms      =  mean( offset_dev_s × 1000 )\n"
         "mean_abs_offset_error_ms  =  mean( |offset_dev_s| × 1000 )",
         "Whether the singer releases notes too early or too late.\n"
         "Only mean and MAE reported (not std); contributes to timing category score.",
         GREEN),
        ("Timing Accuracy  (timing_accuracy)",
         "timing_accuracy  =  mean( |onset_dev_s| × 1000 ≤ τ )   τ = 50 ms",
         "Fraction of notes started within 50 ms of the reference onset.\n"
         "Primary timing accuracy component — weight 0.50 in timing category score.",
         RGBColor(0x1A, 0x5C, 0x96)),
        ("IOI MAE  (ioi_mae_ms)",
         "IOI-MAE  =  mean( |pred_IOI[i] − ref_IOI[i]| × 1000 )\n"
         "where IOI[i] = onset[i+1] − onset[i]  for i = 0 … min(N_pred, N_ref)−2",
         "Inter-onset interval consistency — measures rhythmic stability.\n"
         "Breakpoints: 0→100, 30→88, 60→75, 120→50, 240→0 ms.\n"
         "Combined with onset std-dev into 'rhythm stability' component (weight 0.20).",
         ACCENT),
    ]

    for i, (title, form, desc, clr) in enumerate(items):
        col, row = i % 2, i // 2
        bx, by = 0.25 + col*6.55, 1.18 + row*3.08
        box(sl, bx, by, 6.2, 2.96, fill=RGBColor(0xF0, 0xFB, 0xF0), line=clr)
        box(sl, bx, by, 6.2, 0.36, fill=clr)
        txt(sl, title, bx+0.1, by+0.04, 6.0, 0.3, size=13, bold=True, color=WHITE)
        box(sl, bx+0.1, by+0.42, 6.0, 0.72, fill=RGBColor(0xE0, 0xF4, 0xE0), line=clr)
        txt(sl, form, bx+0.18, by+0.46, 5.84, 0.64, size=11, italic=True, color=DARK_BLUE)
        txt(sl, desc, bx+0.1, by+1.22, 6.0, 1.66, size=11, color=TEXT_DARK)


def slide_duration_lyric_metrics(prs):
    sl = add_slide(prs)
    header(sl, "Metrics — Duration & Lyric",
           "Duration accuracy and phoneme boundary timing · from DurationMetrics and LyricMetrics")

    # Duration (left)
    box(sl, 0.25, 1.15, 6.1, 6.2, fill=RGBColor(0xFF, 0xF8, 0xF0))
    txt(sl, "Duration Metrics", 0.38, 1.22, 5.8, 0.38, size=14, bold=True, color=ACCENT)

    dur = [
        ("Duration Error  (mean_abs_duration_error_s)",
         "error = pred_dur − ref_dur\n"
         "mean_duration_error_s     = mean(error)\n"
         "mean_abs_duration_error_s = mean(|error|)",
         "Absolute difference between predicted and reference note duration.\n"
         "Singer can start correctly but still hold too short or too long."),
        ("Relative Duration Error  (mean_relative_duration_error)",
         "RelErr  =  mean( |pred_dur − ref_dur| / ref_dur )",
         "Tempo-independent: 100 ms error matters more on a 200 ms note\n"
         "than on a 2 s note.  Primary component weight: 0.60.\n"
         "Breakpoints: 0→100, 0.1→90, 0.2→75, 0.5→50, 1.0→0."),
        ("Duration Ratio  (mean_duration_ratio)",
         "DurRatio  =  mean( pred_dur / ref_dur )\n"
         "RatioDev  =  |DurRatio − 1.0|",
         "1.0 = perfect  |  >1.0 = stretched  |  <1.0 = cut short.\n"
         "Component weight: 0.20.  Breakpoints (on |ratio−1|): 0→100, 0.1→90, …"),
    ]
    for i, (name, form, desc) in enumerate(dur):
        by = 1.68 + i*1.85
        box(sl, 0.38, by, 5.82, 1.76, fill=WHITE, line=ACCENT)
        txt(sl, name, 0.48, by+0.06, 5.62, 0.3, size=12, bold=True, color=ACCENT)
        txt(sl, form, 0.48, by+0.4, 5.62, 0.52, size=11, italic=True, color=DARK_BLUE)
        txt(sl, desc, 0.48, by+0.98, 5.62, 0.7, size=10, color=TEXT_DARK)

    # Lyric (right)
    box(sl, 6.65, 1.15, 6.45, 6.2, fill=RGBColor(0xF5, 0xEE, 0xFF))
    txt(sl, "Lyric / Phoneme Metrics", 6.78, 1.22, 6.2, 0.38, size=14, bold=True, color=PURPLE)

    lyr = [
        ("Phoneme Boundary Error  (mean_abs_phoneme_boundary_error_ms)",
         "mae  =  mean( |phoneme_onset_dev_s| × 1000 )\nfor all PhonemeAlignmentMatch pairs",
         "Positive deviation = predicted phoneme starts later than reference.\n"
         "Breakpoints: 0→100, 15→88, 30→75, 60→50, 120→0 ms.\n"
         "Component weight in lyric score: 0.15."),
        ("Phoneme Overlap Accuracy  (phoneme_overlap_accuracy)",
         "overlap_acc  =  fraction where overlap_fraction ≥ 0.5",
         "Whether the predicted phoneme segment significantly overlaps the reference.\n"
         "Component weight: 0.25.  Higher is better."),
        ("Word Alignment Accuracy  (word_alignment_accuracy)",
         "word_acc  =  |matched_words| / |total_ref_words|",
         "Fraction of reference words that matched a predicted word event.\n"
         "Component weight: 0.35 — highest weight in lyric category."),
        ("Label Match Rate  (label_match_rate)",
         "label_match  =  fraction where phoneme_label identical",
         "Correct phoneme label predicted (not just timing).\n"
         "Component weight: 0.25."),
    ]
    for i, (name, form, desc) in enumerate(lyr):
        by = 1.68 + i*1.42
        box(sl, 6.78, by, 6.2, 1.34, fill=WHITE, line=PURPLE)
        txt(sl, name, 6.88, by+0.05, 6.0, 0.28, size=11, bold=True, color=PURPLE)
        txt(sl, form, 6.88, by+0.38, 6.0, 0.36, size=10, italic=True, color=DARK_BLUE)
        txt(sl, desc, 6.88, by+0.78, 6.0, 0.5, size=10, color=TEXT_DARK)


def slide_normalization(prs):
    sl = add_slide(prs)
    header(sl, "Scoring — Normalization",
           "All raw metrics are converted to [0, 100] via piecewise-linear curves tuned to musical thresholds")

    # Three normalization modes from normalization.py
    modes = [
        ("bounded_score(value, lower, upper)",
         "S  =  100 · clip( (upper − value) / (upper − lower),  0, 1 )",
         "Linear clamp.  lower = best (→100), upper = worst (→0).\n"
         "Available but not called in default scoring — accuracy fractions\n"
         "are converted directly as  score = accuracy × 100.",
         MID_BLUE),
        ("piecewise_score(value, breakpoints)",
         "Piecewise-linear interpolation through (x, score) breakpoint pairs",
         "x-values ascending, clamped at both ends.\n"
         "Used for all error-based metrics (MACE, onset MAE, phoneme boundary error…).",
         GREEN),
        ("gaussian_penalty(value, sigma)",
         "S  =  100 · exp( −value² / (2σ²) )",
         "Smooth Gaussian decay from 0.  Available mode but not used\n"
         "in current default scoring configuration.",
         ACCENT),
    ]
    for i, (name, form, desc, clr) in enumerate(modes):
        bx, by = 0.25 + i*4.35, 1.2
        box(sl, bx, by, 4.15, 2.3, fill=WHITE, line=clr)
        box(sl, bx, by, 4.15, 0.34, fill=clr)
        txt(sl, name, bx+0.1, by+0.04, 3.95, 0.28, size=12, bold=True, color=WHITE)
        txt(sl, form, bx+0.1, by+0.4, 3.95, 0.42, size=11, italic=True, color=DARK_BLUE)
        txt(sl, desc, bx+0.1, by+0.9, 3.95, 1.32, size=11, color=TEXT_DARK)

    # Breakpoints table
    txt(sl, "Implemented Piecewise Breakpoints  (from source code)",
        0.25, 3.65, 12.8, 0.38, size=13, bold=True, color=DARK_BLUE)

    headers = ["Metric", "Breakpoints (x → score)", "Notes"]
    rows = [
        ["MACE (cents)", "0→100 · 25→88 · 50→75 · 100→50 · 200→0", "50¢ ≈ quarter-tone · 100¢ = 1 semitone"],
        ["Pitch RMSE (cents)", "0→100 · 25→88 · 50→72 · 100→45 · 200→0", "Stricter than MACE at 50¢+"],
        ["Onset MAE (ms)", "0→100 · 25→88 · 50→75 · 100→50 · 200→0", "50 ms = one timing tolerance window"],
        ["IOI MAE (ms)", "0→100 · 30→88 · 60→75 · 120→50 · 240→0", "Slightly more lenient than onset"],
        ["Onset Std-Dev (ms)", "0→100 · 30→88 · 60→72 · 120→45 · 240→0", "Rhythmic consistency measurement"],
        ["Relative Duration Error", "0→100 · 0.1→90 · 0.2→75 · 0.5→50 · 1.0→0", "Dimensionless; tempo-independent"],
        ["Duration Ratio Dev (|ratio−1|)", "0→100 · 0.1→90 · 0.25→75 · 0.5→50 · 1.0→0", "0 = perfect · 0.5 = 50% off"],
        ["Phoneme Boundary MAE (ms)", "0→100 · 15→88 · 30→75 · 60→50 · 120→0", "Stricter than note timing"],
    ]
    table(sl, headers, rows, 0.25, 4.08, [2.6, 5.5, 4.7],
          row_h=0.36, fsize=11)


def slide_category_scores(prs):
    sl = add_slide(prs)
    header(sl, "Scoring — Category Scores",
           "Four scoring functions · all use confidence-weighted aggregation  (confidence = min(1, n_evaluated/3))")

    cats = [
        ("Pitch Score", MID_BLUE,
         [("Accuracy (pitch_accuracy)", "0.50"),
          ("Intonation — MACE (mace_cents)", "0.30"),
          ("Stability — RMSE (pitch_rmse_cents)", "0.20")],
         "scoring/pitch_scoring.py · compute_pitch_score()"),
        ("Timing Score", GREEN,
         [("Timing Accuracy (timing_accuracy)", "0.50"),
          ("Onset MAE (mean_abs_onset_error_ms)", "0.30"),
          ("Rhythm Stability — IOI + std", "0.20")],
         "scoring/timing_scoring.py · compute_timing_score()"),
        ("Duration Score", ACCENT,
         [("Relative Error (mean_relative_duration_error)", "0.60"),
          ("Ratio Dev (|mean_duration_ratio − 1|)", "0.20"),
          ("Phrase Consistency — std + rel error", "0.20")],
         "scoring/duration_scoring.py · compute_duration_score()"),
        ("Lyric Score", PURPLE,
         [("Word Alignment (word_alignment_accuracy)", "0.35"),
          ("Phoneme Overlap (phoneme_overlap_accuracy)", "0.25"),
          ("Label Match (label_match_rate)", "0.25"),
          ("Boundary Timing (mean_abs_phoneme_boundary_error_ms)", "0.15")],
         "scoring/lyric_scoring.py · compute_lyric_clarity_score()"),
    ]

    for i, (title, clr, components, src) in enumerate(cats):
        col, row = i % 2, i // 2
        bx, by = 0.25 + col*6.55, 1.18 + row*3.1
        box(sl, bx, by, 6.2, 3.0, fill=WHITE, line=clr)
        box(sl, bx, by, 6.2, 0.36, fill=clr)
        txt(sl, title, bx+0.1, by+0.04, 6.0, 0.3, size=14, bold=True, color=WHITE)
        for j, (comp, w) in enumerate(components):
            cy = by + 0.44 + j*0.52
            box(sl, bx+0.12, cy, 5.96, 0.44, fill=LGREY, line=MGREY)
            txt(sl, comp, bx+0.22, cy+0.06, 4.0, 0.3, size=11, color=TEXT_DARK)
            txt(sl, f"w={w}", bx+4.2, cy+0.06, 1.75, 0.3, size=12,
                bold=True, color=clr, align=PP_ALIGN.RIGHT)
        fy = by + 0.44 + len(components)*0.52
        box(sl, bx+0.12, fy+0.06, 5.96, 0.38, fill=RGBColor(0xEE, 0xEE, 0xFF), line=clr)
        txt(sl, src, bx+0.22, fy+0.1, 5.76, 0.28, size=10, italic=True, color=clr)


def slide_final_score(prs):
    sl = add_slide(prs)
    header(sl, "Scoring — Final Score & Interpretation",
           "Confidence-weighted aggregation  ·  deterministic rule-based feedback  ·  no LLMs")

    # Final formula
    box(sl, 0.4, 1.15, 12.5, 1.1, fill=DARK_BLUE)
    txt(sl,
        "S_Final  =  0.40 · S_Pitch  +  0.30 · S_Timing  +  0.15 · S_Duration  +  0.15 · S_Lyric",
        0.55, 1.26, 12.2, 0.55, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "Weights: Pitch (40%)  ·  Timing (30%)  ·  Duration (15%)  ·  Lyric (15%)",
        0.55, 1.78, 12.2, 0.36, size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # Confidence-weighting box
    box(sl, 0.4, 2.24, 12.5, 0.66, fill=RGBColor(0xF0, 0xF5, 0xFF), line=MID_BLUE)
    txt(sl, "Confidence-weighted aggregation  (from performance_scoring.py):",
        0.55, 2.28, 12.2, 0.3, size=12, bold=True, color=DARK_BLUE)
    txt(sl,
        "overall = Σ( category_score · nominal_weight · confidence ) / Σ( nominal_weight · confidence )"
        "   where   confidence = min(1.0, n_evaluated / 3.0)",
        0.55, 2.57, 12.2, 0.28, size=11, italic=True, color=DARK_BLUE)

    # Interpretation table
    txt(sl, "Score Interpretation  (scoring/interpretation.py — rule-based only)",
        0.4, 3.05, 12.5, 0.38, size=13, bold=True, color=DARK_BLUE)

    headers = ["Score Range", "Level", "Coaching Interpretation"]
    rows = [
        ["≥ 90", "excellent", "Performance-ready — only minor refinements needed"],
        ["≥ 75", "good",      "Mostly correct — a few targeted corrections needed"],
        ["≥ 55", "fair",      "Noticeable issues that require focused practice"],
        ["< 55", "needs_work","Focused coaching needed before full-song practice"],
    ]
    table(sl, headers, rows, 0.4, 3.48, [1.6, 1.8, 9.1], row_h=0.48, fsize=13)

    # Feedback
    box(sl, 0.4, 5.54, 12.5, 1.84, fill=LGREY, line=MGREY)
    txt(sl, "Feedback messages (all rule-based, static lookup table):",
        0.55, 5.6, 12.2, 0.34, size=12, bold=True, color=DARK_BLUE)
    blist(sl, [
        "Strengths  (excellent / good): e.g. 'Strong pitch intonation and accuracy'  |  "
        "'Excellent rhythmic timing'  |  'Well-controlled note durations'",
        "Weaknesses (fair / needs_work): e.g. 'Pitch intonation needs significant improvement'  |  "
        "'Significant timing inconsistencies detected'",
    ], 0.55, 5.98, 12.2, fsize=11, gap=0.62)


def slide_pitch_model(prs):
    sl = add_slide(prs)
    header(sl, "Model — Pitch Extraction  (models/pitch/)",
           "torchcrepe (primary) or pYIN (fallback) + WebRTC VAD · outputs f̂ₜ, v̂ₜ at 100 fps")

    # Pipeline flow
    steps = [
        ("WAV Audio  (16 kHz, mono)", DARK_BLUE),
        ("WebRTC VAD  (aggressiveness=2, frame=20 ms, smoothing_window=5)\n"
         "→ voiced_mask at VAD frame rate", RGBColor(0x15, 0x5F, 0x9F)),
        ("PitchModelWrapper.predict()\n"
         "torchcrepe: CREPE-style CNN, Viterbi decoder, periodicity_threshold=0.21\n"
         "pYIN fallback: librosa.pyin, fmin=65 Hz, fmax=2093 Hz\n"
         "→ (times, f0_raw, confidence) at 100 fps (hop=160 @ 16 kHz)", RGBColor(0x0F, 0x4A, 0x7A)),
        ("align_vad_to_pitch()  — resample VAD mask to pitch frame rate", RGBColor(0x0A, 0x3A, 0x60)),
        ("fuse_vad_and_pitch()  — full VAD-pitch fusion and contour cleaning\n"
         "1. voiced_final = AND(pitch_voiced, vad_aligned)  →  2. zero unvoiced frames\n"
         "3. interpolate short gaps (linear, max=10 frames, updates voiced_final)\n"
         "4. median filter voiced frames (kernel=5)  →  5. Gaussian smooth (σ=1.0)", DARK_BLUE),
        ("Output: timestamps (100 fps), f0 (Hz), voiced (bool)", RGBColor(0x05, 0x25, 0x42)),
    ]
    for i, (label, clr) in enumerate(steps):
        by = 1.18 + i*1.02
        box(sl, 0.25, by, 0.6, 0.88, fill=clr)
        txt(sl, str(i+1), 0.28, by+0.22, 0.54, 0.44, size=14, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        box(sl, 0.95, by, 9.5, 0.88, fill=WHITE, line=clr)
        txt(sl, label, 1.05, by+0.06, 9.3, 0.76, size=12, color=TEXT_DARK)
        if i < len(steps)-1:
            txt(sl, "▼", 0.38, by+0.9, 0.5, 0.2, size=10, color=DARK_BLUE,
                align=PP_ALIGN.CENTER)

    # Config summary
    box(sl, 10.6, 1.18, 2.55, 7.15, fill=RGBColor(0xE8, 0xF4, 0xFF), line=MID_BLUE)
    txt(sl, "Config\n(pitch.yaml)", 10.7, 1.24, 2.35, 0.5, size=12,
        bold=True, color=DARK_BLUE)
    cfg_items = [
        "backend: torchcrepe",
        "hop_length: 160",
        "fmin: 50.0 Hz",
        "fmax: 1000.0 Hz",
        "model_capacity: full",
        "use_viterbi: true",
        "periodicity: 0.21",
        "silence_db: -60.0",
        "VAD agg: 2",
        "VAD frame: 20 ms",
    ]
    for j, c in enumerate(cfg_items):
        txt(sl, c, 10.7, 1.78+j*0.55, 2.35, 0.48, size=11, color=TEXT_DARK)


def slide_onset_model(prs):
    sl = add_slide(prs)
    header(sl, "Model — Onset / Offset Detection  (models/onset_offset/model_v5.py)",
           "Wav2Vec2OnsetDetector  ·  facebook/wav2vec2-base backbone  ·  3-class linear head  ·  50 Hz output")

    # Architecture
    arch = [
        ("Input", "Raw 16 kHz waveform  [B, T_samples]",
         "Audio loaded once by pipeline; passed directly as float32 tensor",
         DARK_BLUE),
        ("Feature Extractor", "Wav2Vec2 CNN Feature Extractor  (frozen)",
         "7-layer strided Conv1d stack · total stride = 320 samples → 50 Hz frame rate\n"
         "freeze_feature_extractor=True  (pretrained weights locked)",
         MID_BLUE),
        ("Transformer Encoder", "Wav2Vec2 12-layer Transformer  (fine-tuned)",
         "hidden_size=768 · 12 attention heads · dropout=0.1\n"
         "All transformer layers are updated during fine-tuning",
         RGBColor(0x27, 0x6F, 0x96)),
        ("Dropout + Head", "Dropout(p=0.1) → Linear(768 → 3)",
         "One shared linear head producing 3 logits per frame:\n"
         "logits[:, 0] = onset  ·  logits[:, 1] = offset  ·  logits[:, 2] = active (voiced)",
         RGBColor(0x1A, 0x5C, 0x96)),
        ("Output", "{onset_logits, offset_logits, active_logits}  [B, T_frames]",
         "Apply sigmoid → probabilities · peak-pick → {ŝₙ, êₙ} note boundaries\n"
         "Frame rate: T_frames = T_samples / 320  @  50 Hz",
         RGBColor(0xFF, 0x7F, 0x0E)),
    ]

    for i, (tag, title, detail, clr) in enumerate(arch):
        by = 1.18 + i*1.12
        box(sl, 0.25, by, 1.35, 0.96, fill=clr)
        txt(sl, tag, 0.3, by+0.28, 1.25, 0.44, size=12, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        box(sl, 1.7, by, 11.4, 0.96, fill=WHITE, line=clr)
        txt(sl, title,  1.82, by+0.04, 11.15, 0.36, size=13, bold=True, color=clr)
        txt(sl, detail, 1.82, by+0.46, 11.15, 0.44, size=11, color=TEXT_DARK)
        if i < len(arch)-1:
            txt(sl, "▼", 0.73, by+0.98, 0.6, 0.22, size=11, color=DARK_BLUE,
                align=PP_ALIGN.CENTER)

    # Training / config note
    box(sl, 0.25, 6.85, 12.8, 0.62, fill=DARK_BLUE)
    txt(sl,
        "Training:  fine-tune from facebook/wav2vec2-base  ·  freeze CNN feature extractor  ·  "
        "fine-tune transformer layers  ·  BCE loss on 3 heads  ·  "
        "pipeline.yaml checkpoint:  checkpoints/best.ckpt",
        0.38, 6.88, 12.55, 0.55, size=11, color=WHITE)


def slide_phoneme_model(prs):
    sl = add_slide(prs)
    header(sl, "Model — Phoneme Boundary Detection  (models/phoneme/phoneme_model.py)",
           "facebook/wav2vec2-lv-60-espeak-cv-ft  ·  CTC alignment  ·  post-processing  ·  50 Hz")

    # Pipeline steps
    steps = [
        ("WAV Audio  (16 kHz, mono)",
         PURPLE),
        ("Wav2Vec2FeatureExtractor  —  CNN preprocessing  (pretrained, eval mode)",
         RGBColor(0x5A, 0x2A, 0x82)),
        ("Wav2Vec2 Transformer Encoder  —  contextual audio embeddings\n"
         "model: facebook/wav2vec2-lv-60-espeak-cv-ft  (60-language multilingual)\n"
         "tokenizer: Wav2Vec2PhonemeCTCTokenizer  (eSpeak IPA vocabulary)",
         RGBColor(0x49, 0x1F, 0x6E)),
        ("CTC Argmax Alignment  →  frame-level phoneme token predictions\n"
         "collapse_repeated_tokens=True · remove_blank_tokens=True",
         RGBColor(0x39, 0x14, 0x5A)),
        ("Post-processing  (two passes):\n"
         "1. Split long segments: max_segment_ms=300, min_split_prob=0.25\n"
         "2. Blank-region scan: insert phonemes missed by CTC blank dominance (min_phoneme_prob=0.05)",
         RGBColor(0x29, 0x0A, 0x47)),
        ("Output: List[PhonemeSegment(phoneme, start_time, end_time, confidence)]",
         RGBColor(0x1A, 0x03, 0x33)),
    ]
    for i, (label, clr) in enumerate(steps):
        by = 1.18 + i*1.0
        box(sl, 0.25, by, 0.6, 0.86, fill=clr)
        txt(sl, str(i+1), 0.28, by+0.21, 0.54, 0.44, size=14, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        box(sl, 0.95, by, 12.1, 0.86, fill=WHITE, line=clr)
        txt(sl, label, 1.05, by+0.05, 11.9, 0.76, size=12, color=TEXT_DARK)
        if i < len(steps)-1:
            txt(sl, "▼", 0.38, by+0.88, 0.5, 0.18, size=10, color=DARK_BLUE,
                align=PP_ALIGN.CENTER)

    formula(sl,
            "Scored via lyric_metrics.py:  "
            "mean_abs_phoneme_boundary_error_ms  =  mean( |onset_deviation_s| × 1000 )"
            "  over all PhonemeAlignmentMatch pairs  (onset deviation only, not end-boundary)",
            0.25, 7.15, 12.8, 0.38)


def slide_results_pitch(prs):
    sl = add_slide(prs)
    header(sl, "Results — Pitch Model  (NanoPitch, Kim submission)",
           "torchcrepe backend · staged curriculum training · 100 clips per SNR condition · GTSinger")

    txt(sl, "Raw Pitch Accuracy (RPA) and Voicing Detection Rate across SNR conditions",
        0.25, 1.22, 12.8, 0.38, size=13, bold=True, color=DARK_BLUE)

    headers = ["Condition", "VAD Acc", "RPA (offline)", "RCA (offline)",
               "Gross Err", "Median Cents", "RPA (realtime)", "Med Cents (RT)"]
    rows = [
        ["-5 dB",  "93.8%", "90.7%", "91.3%", "9.3%",  "56.6 ¢", "89.2%", "53.6 ¢"],
        ["+0 dB",  "93.9%", "89.4%", "89.4%", "10.7%", "56.2 ¢", "87.0%", "57.2 ¢"],
        ["+5 dB",  "92.7%", "90.3%", "90.5%", "9.75%", "33.7 ¢", "88.5%", "35.4 ¢"],
        ["+10 dB", "96.2%", "94.2%", "94.3%", "5.8%",  "15.7 ¢", "91.6%", "17.1 ¢"],
        ["+20 dB", "96.6%", "93.2%", "93.2%", "6.8%",  "14.9 ¢", "90.1%", "31.5 ¢"],
        ["Clean",  "97.9%", "92.5%", "92.8%", "7.5%",  "14.7 ¢", "90.8%", "15.8 ¢"],
    ]
    table(sl, headers, rows, 0.25, 1.65,
          [1.6, 1.1, 1.55, 1.55, 1.2, 1.7, 1.6, 2.0],
          row_h=0.42, fsize=12)

    box(sl, 0.25, 4.76, 12.8, 0.06, fill=ACCENT)
    txt(sl, "Key observations", 0.25, 4.88, 12.8, 0.34, size=13, bold=True, color=DARK_BLUE)

    obs = [
        "Best offline RPA: 94.2% at +10 dB SNR — within MIREX competition-level accuracy",
        "Clean audio: median cent error 14.7 ¢ offline — well below the ±50 ¢ coaching tolerance threshold",
        "Noisy conditions (−5 dB): offline RPA still above 90% — robust to moderate background noise",
        "Realtime mode trades ~2–3% RPA for lower latency (suitable for streaming / live feedback)",
        "VAD accuracy ≥ 92.7% across all conditions — reliable voiced/unvoiced gating",
    ]
    for i, o in enumerate(obs):
        txt(sl, "  •  " + o, 0.25, 5.26 + i*0.38, 12.8, 0.34, size=12, color=TEXT_DARK)


def slide_results_onset(prs):
    sl = add_slide(prs)
    header(sl, "Results — Note Onset / Offset Detection  (Wav2Vec2 model)",
           "Wav2Vec2OnsetDetector · facebook/wav2vec2-base · GTSinger · 50 ms tolerance")

    # What the model predicts and how it feeds the scoring pipeline
    txt(sl, "Model outputs and their role in the scoring pipeline",
        0.25, 1.22, 12.8, 0.38, size=13, bold=True, color=DARK_BLUE)

    headers = ["Output", "How it is produced", "Used to compute"]
    rows = [
        ["onset_logits  [B, T]",
         "Linear(768→3)[:,0] from Wav2Vec2 last_hidden_state\n→ sigmoid → peak-pick → ŝₙ",
         "OnsetError · TimingAccuracy · IOI-MAE · DurationError"],
        ["offset_logits  [B, T]",
         "Linear(768→3)[:,1] from Wav2Vec2 last_hidden_state\n→ sigmoid → peak-pick → êₙ",
         "OffsetError · DurationError · RelativeDurationError · DurationRatio"],
        ["active_logits  [B, T]",
         "Linear(768→3)[:,2] from Wav2Vec2 last_hidden_state\n→ sigmoid → voiced activity mask",
         "Supplements VAD for voiced-frame gating during note event construction"],
    ]
    table(sl, headers, rows, 0.25, 1.65, [2.8, 5.3, 4.7], row_h=0.88, fsize=12)

    # Peak-picking post-processing
    txt(sl, "Peak-picking post-processing  (shared with all backends)",
        0.25, 4.42, 12.8, 0.36, size=13, bold=True, color=DARK_BLUE)

    pp_items = [
        ("onset_threshold = 0.3", "Minimum sigmoid probability to register a peak as an onset"),
        ("offset_threshold = 0.3", "Minimum sigmoid probability to register a peak as an offset"),
        ("min_distance_frames = 3", "Minimum 3 frames (~60 ms at 50 Hz) between consecutive peaks of the same type"),
        ("pair_onsets_offsets()", "Greedily pair each onset with the nearest subsequent offset → note event list"),
    ]
    for i, (param, desc) in enumerate(pp_items):
        by = 4.84 + i*0.56
        box(sl, 0.25, by, 3.2, 0.48, fill=RGBColor(0xE8, 0xF1, 0xFB), line=MID_BLUE)
        txt(sl, param, 0.35, by+0.08, 3.0, 0.32, size=12, bold=True, color=MID_BLUE)
        txt(sl, desc, 3.55, by+0.08, 9.45, 0.32, size=12, color=TEXT_DARK)

    # Frame rate note
    box(sl, 0.25, 7.1, 12.8, 0.42, fill=DARK_BLUE)
    txt(sl,
        "Frame rate: Wav2Vec2 CNN stride = 320 samples @ 16 kHz  →  50 Hz  (20 ms / frame)  ·  "
        "output resampled to 100 fps canonical grid via fusion/alignment.py",
        0.38, 7.14, 12.5, 0.34, size=11, color=WHITE)


def slide_results_note_pitch(prs):
    sl = add_slide(prs)
    header(sl, "Results — Note-Level Pitch Accuracy  (pilot evaluation)",
           "7-note excerpt · torchcrepe backend · τ = 50¢ · PitchModel_Kim/note_pitch_results.json")

    box(sl, 0.25, 1.15, 12.8, 0.52, fill=DARK_BLUE)
    txt(sl,
        "NotePitchAcc₅₀  =  85.7%   (6 / 7 notes correct)     "
        "Song: 'North Wind Meets the Sea'   Notes: C4 – D4 – E4 – Bb3 – C4 – A3 – D4",
        0.4, 1.22, 12.5, 0.38, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    headers = ["Note", "Lyric", "Ref MIDI", "Ref Hz",
               "Detected MIDI", "Detected Hz", "Median Cents", "Dir", "Correct"]
    rows = [
        ["1", "where", "C4 (60.0)",  "261.6", "59.74",  "257.7", "−26.1 ¢", "Flat",  "✓"],
        ["2", "the",   "D4 (62.0)",  "293.7", "60.47",  "268.9", "−152.8 ¢","Flat",  "✗"],
        ["3", "north", "E4 (64.0)",  "329.6", "64.08",  "331.2", "+8.2 ¢",  "Sharp", "✓"],
        ["4", "wind",  "Bb3 (58.0)", "233.1", "57.63",  "228.2", "−36.8 ¢", "Flat",  "✓"],
        ["5", "meets", "C4 (60.0)",  "261.6", "60.32",  "266.6", "+32.5 ¢", "Sharp", "✓"],
        ["6", "the",   "A3 (57.0)",  "220.0", "57.44",  "225.7", "+44.4 ¢", "Sharp", "✓"],
        ["7", "sea",   "D4 (62.0)",  "293.7", "62.24",  "297.7", "+23.7 ¢", "Sharp", "✓"],
    ]
    table(sl, headers, rows, 0.25, 1.73,
          [0.6, 0.85, 1.35, 0.95, 1.55, 1.45, 1.6, 1.1, 1.3],
          row_h=0.42, fsize=11)

    box(sl, 0.25, 4.86, 12.8, 0.05, fill=ACCENT)
    txt(sl, "Observations", 0.25, 4.97, 12.8, 0.34, size=13, bold=True, color=DARK_BLUE)
    obs = [
        "Note 2 ('the' – D4): only failure · −152.8 ¢ flat (≈1.5 semitones) on a very short 178 ms note "
        "with only 18 voiced frames — insufficient data for reliable median",
        "All 6 correct notes are within ±45 ¢ — Note 6 at 44.4 ¢ is the closest to the ±50 ¢ tolerance boundary",
        "Mild systematic bias: flat on longer sustained notes (Notes 1, 4) and sharp on short/fast notes",
        "Total voiced frames scored: 522 frames across 7 notes — strong voicing detection on this excerpt",
    ]
    for i, o in enumerate(obs):
        txt(sl, "  •  " + o, 0.25, 5.36 + i*0.44, 12.8, 0.38, size=12, color=TEXT_DARK)


def slide_pipeline_summary(prs):
    sl = add_slide(prs)
    header(sl, "Integrated System — VocalCoach_Kim  (7-phase pipeline)",
           "inference/pipeline.py · UnifiedInferencePipeline.predict(wav, musicxml, textgrid, compute_scores=True)")

    phases = [
        ("1", "Audio Preprocessing",
         "utils/audio.py  |  16 kHz mono, normalize  |  "
         "canonical 10 ms frame grid (hop=160)",
         DARK_BLUE),
        ("2", "Pitch + VAD",
         "models/pitch/pipeline.py  |  "
         "torchcrepe → f0 (Hz), voiced mask  @  100 fps",
         MID_BLUE),
        ("3", "Phoneme Extraction",
         "models/phoneme/phoneme_model.py  |  "
         "Wav2Vec2 + CTC  →  List[PhonemeSegment]",
         PURPLE),
        ("4", "Onset / Offset Detection",
         "models/onset_offset/model_v5.py  |  "
         "Wav2Vec2 → {onset, offset, active} logits  →  peak-pick  →  ŝₙ, êₙ",
         ACCENT),
        ("5", "Feature Fusion",
         "fusion/  |  merge_model_outputs → FrameAlignedFeatures (100fps)  |  "
         "build_note_events, build_lyric_events  →  FusedPerformanceRepresentation",
         RGBColor(0x27, 0x6F, 0x96)),
        ("6", "Reference Parsing + Alignment",
         "reference/musicxml_parser + textgrid_parser  →  ReferencePerformanceRepresentation  |  "
         "alignment/reference_alignment  →  AlignmentResult",
         GREEN),
        ("7", "Metrics + Scoring",
         "metrics/  →  PerformanceMetricsReport  |  "
         "scoring/  →  PerformanceScoreReport  |  "
         "scoring/interpretation  →  InterpretationSummary",
         RGBColor(0x8C, 0x56, 0x4B)),
    ]
    for i, (num, title, detail, clr) in enumerate(phases):
        by = 1.18 + i*0.88
        box(sl, 0.25, by, 0.72, 0.78, fill=clr)
        txt(sl, num, 0.28, by+0.18, 0.66, 0.44, size=16, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        box(sl, 1.07, by, 12.03, 0.78, fill=WHITE, line=clr)
        txt(sl, title,  1.18, by+0.04, 3.2, 0.34, size=13, bold=True, color=clr)
        txt(sl, detail, 1.18, by+0.42, 11.8, 0.28, size=11, color=TEXT_DARK)
        if i < len(phases)-1:
            txt(sl, "▼", 0.44, by+0.8, 0.5, 0.16, size=9, color=DARK_BLUE,
                align=PP_ALIGN.CENTER)

    box(sl, 0.25, 7.35, 12.8, 0.54, fill=DARK_BLUE)
    txt(sl,
        "result = pipeline.predict(wav, musicxml_path=..., textgrid_path=..., "
        "compute_metrics=True, compute_scores=True)  "
        "→  result.scores.overall_score  ·  result.interpretation.strengths  ·  result.scores.pitch_score  …",
        0.38, 7.38, 12.5, 0.48, size=11, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# BUILD
# ══════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()

    slide_title(prs)              # 1
    slide_overview(prs)           # 2
    slide_parameters(prs)         # 3  — parameters
    slide_pitch_metrics(prs)      # 4  — metrics: pitch
    slide_timing_metrics(prs)     # 5  — metrics: timing
    slide_duration_lyric_metrics(prs)  # 6 — metrics: duration + lyric
    slide_normalization(prs)      # 7  — scoring: normalization
    slide_category_scores(prs)    # 8  — scoring: 4 categories
    slide_final_score(prs)        # 9  — scoring: final + interpretation
    slide_pitch_model(prs)        # 10 — model: pitch
    slide_onset_model(prs)        # 11 — model: Wav2Vec2 onset/offset
    slide_phoneme_model(prs)      # 12 — model: phoneme
    slide_results_pitch(prs)      # 13 — results: pitch
    slide_results_onset(prs)      # 14 — results: onset/offset
    slide_results_note_pitch(prs) # 15 — results: note-level pitch
    slide_pipeline_summary(prs)   # 16 — integrated pipeline

    out = ("/mnt/researchfiles/ECE IMAPLE/cluster_data/user_data/pc833"
           "/Voice-Coaching-System/VoiceCoachingSystem_Presentation_v2.pptx")
    prs.save(out)
    print(f"Saved  →  {out}")
    print(f"Slides:    {len(prs.slides)}")


if __name__ == "__main__":
    build()
